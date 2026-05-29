import os
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PyPDF2 import PdfReader


def parse_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    parsers = {
        ".docx": _parse_docx,
        ".pptx": _parse_pptx,
        ".xlsx": _parse_xlsx,
        ".xls": _parse_xlsx,
        ".pdf": _parse_pdf,
        ".md": _parse_text,
        ".txt": _parse_text,
        ".csv": _parse_text,
    }
    parser = parsers.get(ext, _parse_text)
    try:
        return parser(file_path)
    except Exception as e:
        return f"[文件解析失败] {str(e)}"


def _parse_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def _parse_pptx(path: str) -> str:
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(paragraph.text.strip())
    return "\n".join(texts)


def _parse_xlsx(path: str) -> str:
    wb = load_workbook(path, data_only=True)
    texts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        texts.append(f"### Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join([str(c) if c is not None else "" for c in row])
            if row_text.strip():
                texts.append(row_text)
    return "\n".join(texts)


def _parse_pdf(path: str) -> str:
    reader = PdfReader(path)
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text)
    return "\n".join(texts)


def _parse_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()
